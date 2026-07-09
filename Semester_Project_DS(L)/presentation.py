from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(26, 35, 126)
GOLD = RGBColor(255, 215, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(66, 165, 245)
GREEN = RGBColor(76, 175, 80)
RED = RGBColor(244, 67, 54)
ORANGE = RGBColor(255, 152, 0)
PURPLE = RGBColor(156, 39, 176)
TEAL = RGBColor(0, 150, 136)
DARK_GRAY = RGBColor(33, 33, 33)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide():
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.color.rgb = DARK_BLUE
    
    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(1.5), prs.slide_width, Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.color.rgb = GOLD
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.3), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "EMAIL / SMS SPAM CLASSIFIER"
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.3), Inches(1))
    sub_frame = sub_box.text_frame
    sub_frame.text = "Project Report 2026"
    sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sub_frame.paragraphs[0].font.size = Pt(28)
    sub_frame.paragraphs[0].font.color.rgb = GOLD
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Section: E | Department of Computer Science | UET Peshawar\nInstructor: Rida Zarkash"
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_frame.paragraphs[0].font.size = Pt(16)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

def add_content_slide(title, content_lines, bullet_points=None, highlight_text=None):
    """Generic content slide template"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = WHITE
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    cf = content_box.text_frame
    cf.text = content_lines
    cf.paragraphs[0].font.size = Pt(20)
    cf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    return slide

def add_slide_with_box(title, boxes):
    """Slide with multiple colored boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = WHITE
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    
    # Add boxes
    colors = [GOLD, LIGHT_BLUE, GREEN, ORANGE, PURPLE, TEAL]
    for i, box_data in enumerate(boxes):
        x = Inches(0.8 + (i % 3) * 3.8)
        y = Inches(1.8 + (i // 3) * 2.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.5), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.color.rgb = colors[i % len(colors)]
        
        text_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(3.3), Inches(2))
        tf = text_box.text_frame
        tf.text = box_data
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_agenda_slide():
    """Slide 3: Agenda"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.text = "📋 TABLE OF CONTENTS"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Left column
    left_items = [
        "🔹 Abstract",
        "🔹 Introduction",
        "🔹 Problem Statement",
        "🔹 Objectives",
        "🔹 Methodology",
        "    ▪ Dataset",
        "    ▪ Data Cleaning",
        "    ▪ Exploratory Data Analysis",
        "    ▪ Text Preprocessing",
        "    ▪ Feature Extraction"
    ]
    
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(5))
    lf = left_box.text_frame
    lf.text = "\n".join(left_items)
    for p in lf.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)
    
    # Right column
    right_items = [
        "🔹 Model Building",
        "🔹 Model Evaluation",
        "🔹 Model Selection",
        "🔹 Advanced Model Development",
        "🔹 Model Saving",
        "🔹 Deployment",
        "🔹 Conclusion"
    ]
    
    right_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.5), Inches(5))
    rf = right_box.text_frame
    rf.text = "\n".join(right_items)
    for p in rf.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

def create_all_slides():
    """Create all presentation slides"""
    
    # Slide 1: Title
    add_title_slide()
    
    # Slide 2: Group Members
    add_slide_with_box(
        "👥 GROUP MEMBERS",
        [
            "01\nHazrat Bilal\n24PWBCS1345",
            "02\nHaroon Ur Rashid\n24PWBCS1357",
            "03\nKhalida Afghan\n24PWBCS1416",
            "04\nIlham Raza\n24PWBCS1259"
        ]
    )
    
    # Slide 3: Agenda
    add_agenda_slide()
    
    # Slide 4: Abstract
    add_content_slide(
        "📄 ABSTRACT",
        """An SMS Spam Detection system using Machine Learning and Natural Language Processing (NLP).

🎯 Objective: Automatically classify SMS messages as Spam or Ham

⚙️ Approach:
• Data Cleaning & Text Preprocessing
• TF-IDF Vectorization
• Naive Bayes Classification

✅ Best Model: Multinomial Naive Bayes (100% Precision)

🚀 Deployment: Streamlit Web Application"""
    )
    
    # Slide 5: Introduction
    add_content_slide(
        "🔍 INTRODUCTION",
        """What are Spam Messages?
📱 Spam = Unsolicited messages containing:
   • Advertisements
   • Scams & Fraudulent offers
   • Phishing attempts

💡 Solution: Machine Learning + NLP

1️⃣ Clean & Preprocess Text
2️⃣ Convert to Numerical Features (TF-IDF)
3️⃣ Train Naive Bayes Models
4️⃣ Deploy as Web Application"""
    )
    
    # Slide 6: Problem Statement
    add_content_slide(
        "⚠️ PROBLEM STATEMENT",
        """📈 Increasing Number of Spam SMS

Class Distribution:
• Ham (Non-Spam): 87.37%
• Spam: 12.63%

🔴 Challenges:
• Users exposed to fraud & phishing
• Manual filtering is time-consuming
• Need for automated, accurate classification"""
    )
    
    # Slide 7: Objectives
    add_content_slide(
        "🎯 OBJECTIVES",
        """1️⃣ Build an SMS Spam Detection Model
2️⃣ Clean and Preprocess Text Data
3️⃣ Perform Exploratory Data Analysis
4️⃣ Convert Text into TF-IDF Features
5️⃣ Train and Compare Naive Bayes Classifiers
6️⃣ Evaluate Model Performance
7️⃣ Deploy the Model using Streamlit"""
    )
    
    # Slide 8: Dataset
    add_content_slide(
        "📊 DATASET",
        """📁 SMS Spam Collection Dataset

Dataset Statistics:
• Total Records: 5,572
• Ham Messages: 4,827 (87.37%)
• Spam Messages: 745 (12.63%)

Sample Data:
┌────────┬──────────────────────────────────────────┐
│ Label  │ Message                                  │
├────────┼──────────────────────────────────────────┤
│ ham    │ Go until jurong point, crazy..           │
│ spam   │ Free entry in 2 a wkly comp...           │
│ ham    │ U dun say so early hor...                │
└────────┴──────────────────────────────────────────┘"""
    )
    
    # Slide 9: Data Cleaning
    add_content_slide(
        "🧹 DATA CLEANING",
        """🔧 Cleaning Steps:

✓ Remove unnecessary columns (Unnamed: 2,3,4)
✓ Rename columns: v1→target, v2→text
✓ Label Encoding: Ham=0, Spam=1
✓ Remove duplicate records
✓ No missing values found

Before:
┌────────┬─────────────┬────────────┬───────────┬──────────┐
│ v1     │ v2          │ Unnamed: 2 │ Unnamed:3 │ Unnamed:4│
├────────┼─────────────┼────────────┼───────────┼──────────┤
│ ham    │ message...  │ NaN        │ NaN       │ NaN      │
└────────┴─────────────┴────────────┴───────────┴──────────┘

After:
┌────────┬──────────────┐
│ target │ text         │
├────────┼──────────────┤
│ 0      │ message...   │
└────────┴──────────────┘"""
    )
    
    # Slide 10: EDA - Class Distribution
    add_content_slide(
        "📊 EDA - CLASS DISTRIBUTION",
        """Class Distribution of SMS Messages

          📊 87.37%              📊 12.63%
         Ham (Non-Spam)            Spam

⚠️ Dataset is Imbalanced!
• Precision is more important than accuracy"""
    )
    
    # Slide 11: EDA - Text Length
    add_content_slide(
        "📊 EDA - TEXT LENGTH",
        """Three New Numerical Features Created:

📏 Characters    📝 Words    📝 Sentences

Example: "You've won a free prize!"
Characters: 22    Words: 5    Sentences: 1

💡 Spam messages are generally:
   • Longer (more characters)
   • More words
   • More sentences"""
    )
    
    # Slide 12: Character Distribution
    add_content_slide(
        "📊 EDA - CHARACTER DISTRIBUTION",
        """Comparison of Character Count

[Insert Histogram Image: media/image14.png]

📌 Key Insight:
• Ham: Most contain < 100 characters
• Spam: Concentrated between 100-170 characters

🔵 Ham    🟠 Spam"""
    )
    
    # Slide 13: Word Distribution
    add_content_slide(
        "📊 EDA - WORD DISTRIBUTION",
        """Comparison of Word Count

[Insert Histogram Image: media/image15.png]

📌 Key Insight:
• Ham: Mostly 5-20 words
• Spam: Concentrated between 20-35 words

🔵 Ham    🟠 Spam"""
    )
    
    # Slide 14: Pair Plot
    add_content_slide(
        "📊 EDA - PAIR PLOT",
        """[Insert Pair Plot Image: media/image16.png]

📌 Key Insights:
• Positive correlation: more characters = more words
• Spam messages show higher values across all features
• Features are useful for classification

🔵 Ham Messages    🟠 Spam Messages"""
    )
    
    # Slide 15: Correlation Heatmap
    add_content_slide(
        "📊 EDA - CORRELATION HEATMAP",
        """[Insert Correlation Heatmap Image: media/image17.png]

📌 Key Insights:
• Strong correlation: chars ↔ words (0.97)
• Moderate: sentences ↔ chars (0.62) & words (0.68)
• Weak positive correlation with target"""
    )
    
    # Slide 16: Text Preprocessing
    add_content_slide(
        "🔧 TEXT PREPROCESSING",
        """Natural Language Processing (NLP) Pipeline

Raw Text: "Free entry in 2 a wkly comp!"
     ↓ 1️⃣ Convert to Lowercase
"free entry in 2 a wkly comp!"
     ↓ 2️⃣ Tokenization
['free', 'entry', 'in', '2', 'a', 'wkly', 'comp!']
     ↓ 3️⃣ Remove Special Characters
['free', 'entry', 'in', '2', 'a', 'wkly', 'comp']
     ↓ 4️⃣ Remove Stop Words & Punctuation
['free', 'entry', 'wkly', 'comp']
     ↓ 5️⃣ Porter Stemming
['free', 'entri', 'wkli', 'comp']"""
    )
    
    # Slide 17: Word Clouds
    add_content_slide(
        "☁️ WORD CLOUDS",
        """🟠 Spam Message Word Cloud    🔵 Ham Message Word Cloud

[Insert Word Cloud Images: media/image19.png & media/image20.png]

📌 Key Differences:
• Spam: Promotional words (free, call, text, claim)
• Ham: Conversational words (love, know, go, good)"""
    )
    
    # Slide 18: Top Words
    add_content_slide(
        "📊 TOP FREQUENT WORDS",
        """🟠 Top 30 Words in Spam Messages

[Insert Bar Chart: media/image21.png]

🔵 Top 30 Words in Ham Messages

[Insert Bar Chart: media/image22.png]"""
    )
    
    # Slide 19: Model Building
    add_content_slide(
        "🤖 MODEL BUILDING",
        """Workflow:
TF-IDF Features → Train/Test Split (80/20) → Naive Bayes Classifiers

🔹 Gaussian Naive Bayes (GNB)
🔹 Multinomial Naive Bayes (MNB)
🔹 Bernoulli Naive Bayes (BNB)

Evaluation Metrics:
• Accuracy
• Precision (most important for imbalanced data)
• Confusion Matrix"""
    )
    
    # Slide 20: Model Evaluation
    add_content_slide(
        "📊 MODEL EVALUATION - NAIVE BAYES",
        """📊 Results Comparison

┌────────────────┬──────────────┬──────────────┐
│ Model          │   Accuracy   │   Precision  │
├────────────────┼──────────────┼──────────────┤
│ Gaussian NB    │   87.33%     │   51.60%     │
│ Multinomial NB │   97.10%     │  🏆100.00%  │
│ Bernoulli NB   │   98.36%     │   99.19%     │
└────────────────┴──────────────┴──────────────┘

🏆 Best Model: Multinomial Naive Bayes
✅ Why? 100% Precision - No ham misclassified as spam

⚠️ Note: Precision > Accuracy for imbalanced datasets"""
    )
    
    # Slide 21: Advanced Models
    add_content_slide(
        "🚀 ADVANCED MODELS",
        """🔬 Additional Models Tested:
🧠 Support Vector Machine (SVM)
🌲 Random Forest (RF)

📊 Advanced Model Results
┌────────────────┬──────────────┬──────────────┐
│ Model          │   Accuracy   │   Precision  │
├────────────────┼──────────────┼──────────────┤
│ SVM            │   97.58%     │   97.48%     │
│ Random Forest  │   97.68%     │   97.50%     │
└────────────────┴──────────────┴──────────────┘

💡 Decision: MNB remains the final model
📌 Reason: 100% Precision > Small accuracy improvement"""
    )
    
    # Slide 22: Model Saving
    add_content_slide(
        "💾 MODEL SAVING",
        """Saved Models using Pickle Library

📁 Saved Files:
  🗂️ model.pkl        (Multinomial NB)
  🗂️ vectorizer.pkl   (TF-IDF Vectorizer)

🔄 These files are used for deployment"""
    )
    
    # Slide 23: Deployment
    add_content_slide(
        "🌐 DEPLOYMENT - STREAMLIT",
        """[Insert Streamlit App Screenshot: media/image23.png]

┌─────────────────────────────────────────────┐
│  📱 SMS / Email Spam Classifier             │
│                                             │
│  ┌─────────────────────────────────────┐   │
│  │ Enter your message here...          │   │
│  └─────────────────────────────────────┘   │
│                                             │
│  [ 🔍 Predict Button ]                     │
│                                             │
│  ┌─────────────────────────────────────┐   │
│  │ ✅ Prediction: Not Spam (Ham)      │   │
│  └─────────────────────────────────────┘   │
└─────────────────────────────────────────────┘

🔹 Real-time classification
🔹 User-friendly interface
🔹 Instant results"""
    )
    
    # Slide 24: Deployment Flow
    add_content_slide(
        "🔄 DEPLOYMENT WORKFLOW",
        """1️⃣ User enters message in text box
              ↓
2️⃣ Click "Predict" button
              ↓
3️⃣ Text is preprocessed (lowercase, tokenize, etc.)
              ↓
4️⃣ Converted to TF-IDF features
              ↓
5️⃣ Trained model predicts class
              ↓
6️⃣ Result displayed: "Spam" or "Not Spam (Ham)"

📱 Quick & Easy!"""
    )
    
    # Slide 25: Results Summary
    add_content_slide(
        "📊 RESULTS SUMMARY",
        """🏆 Best Model: Multinomial Naive Bayes

✅ Accuracy: 97.10%
✅ Precision: 100.00%

📈 Key Achievements:
✓ No ham messages misclassified as spam
✓ High accuracy on imbalanced dataset
✓ Real-time prediction capability
✓ User-friendly web interface"""
    )
    
    # Slide 26: Conclusion
    add_content_slide(
        "✅ CONCLUSION",
        """🎯 Project Successfully Achieved:

1️⃣ Complete end-to-end data science pipeline
2️⃣ 100% Precision → No false positives
3️⃣ Deployed as Streamlit web application
4️⃣ Effective solution for spam message detection

🔄 Complete Data Science Pipeline:
Data → Clean → Analyze → Model → Evaluate → Deploy"""
    )
    
    # Slide 27: Thank You
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.color.rgb = DARK_BLUE
    
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(2))
    tf = thanks.text_frame
    tf.text = "🙏 THANK YOU"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.3), Inches(1))
    sf = subtitle.text_frame
    sf.text = "Email/SMS Spam Classifier"
    sf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sf.paragraphs[0].font.size = Pt(24)
    sf.paragraphs[0].font.color.rgb = WHITE
    
    team = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.5))
    tf2 = team.text_frame
    tf2.text = "Team: Hazrat Bilal • Haroon Ur Rashid • Khalida Afghan • Ilham Raza"
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.8))
    ff = footer.text_frame
    ff.text = "UET Peshawar | Department of Computer Science"
    ff.paragraphs[0].alignment = PP_ALIGN.CENTER
    ff.paragraphs[0].font.size = Pt(16)
    ff.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

# Generate the presentation
if __name__ == "__main__":
    create_all_slides()
    
    # Save the presentation
    prs.save("SMS_Spam_Classifier_Presentation.pptx")
    print("✅ Presentation created successfully!")
    print("📁 File: SMS_Spam_Classifier_Presentation.pptx")
    